from flask import Flask, request, send_file
from flask_cors import CORS, cross_origin

import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO

from dotenv import dotenv_values
import google.generativeai as genai

from utils import extract_json_from_text

config = dotenv_values("./.env")

genai.configure(api_key=config["APIKEY"])
model = genai.GenerativeModel(config['MODELNAME'])

history = []

chat = model.start_chat(history=history)

app = Flask(__name__)
cors = CORS(app, resources={r"/*": {"origins": "*"}})
app.config['CORS_HEADERS'] = 'Content-Type'

@app.route("/create_ppt", methods=["POST"])
def queryLLM():
    if request.method != "POST":
        return "Invalid Request Method"
    
    prompt = f"""
    Using the prompt {request.form['title']} to create a presentation with {request.form['numSlides']} slides and include relevant one image keyword with appropriate bullet points, include no other text except for the json.
    Also make sure to keep these instructions in mind {request.form['additionalInstructions']}
    Using this JSON schema:
        Slide = {{id:int, title:str, points:list(str), image_keywords:list(str)}}
    Return a `list[Slide]`
    """
    res = chat.send_message(prompt)
    data = extract_json_from_text(res.text)

    
    for slide in data:
        for keyword in slide['image_keywords']:
            res=requests.get(f"https://api.unsplash.com//search/photos?client_id={config['CLIENTID']}&query={keyword}&page=1&per_page=1&orientation=squarish")
            if res.status_code == 200:
                slide["image_url"]=(res.json()["results"][0]["urls"]["small"])
                break
            else:
                print(res.text)
                
    return data

@app.route('/download', methods=['POST'])
def create_pptx():
    data = request.get_json()
    slides_data = data.get('slides', [])

    prs = Presentation()

    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
        slide = prs.slides.add_slide(slide_layout)

        # Set the title
        title_shape = slide.shapes.title
        title_shape.text = slide_data['title']
        title_shape.text_frame.paragraphs[0].font.size = Pt(24)
        title_shape.text_frame.paragraphs[0].font.bold = True

        # Add bullet points for each point
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(5.5)
        height = Inches(4)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        for point in slide_data['points']:
            p = text_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.level = 0
            p.space_after = Pt(10)
            p.alignment = PP_ALIGN.LEFT
            p.font.color.rgb = RGBColor(0, 0, 0)

        # Download and add image if available
        if slide_data.get('image_url'):
            try:
                response = requests.get(slide_data['image_url'])

                image_stream = BytesIO(response.content)
                image_height = Inches(3)
                image_left = Inches(6.5)
                image_top = Inches(2)
                slide.shapes.add_picture(image_stream, image_left, image_top, height=image_height)
            except Exception as e:
                print(f"Failed to download or add image: {e}")

    # Save the presentation to a BytesIO object
    output = BytesIO()
    prs.save(output)
    output.seek(0)

    return send_file(output, as_attachment=True, download_name='presentation.pptx', mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
            

@app.route("/rewrite", methods=["POST"])
def rewrite():
    pass

